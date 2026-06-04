from __future__ import annotations

import csv
import html
import json
import mimetypes
import re
from pathlib import Path
from typing import Iterable


TEXT_EXTENSIONS = {
    ".md",
    ".mdx",
    ".txt",
    ".rtf",
    ".xml",
    ".yaml",
    ".yml",
    ".json",
    ".jsonl",
    ".csv",
    ".tsv",
    ".ndjson",
    ".py",
    ".js",
    ".ts",
    ".jsx",
    ".tsx",
    ".rs",
    ".go",
    ".java",
    ".c",
    ".cpp",
    ".h",
    ".rb",
    ".php",
    ".swift",
    ".sql",
    ".sh",
}


class ExtractionError(RuntimeError):
    pass


def extract_text(path: Path, max_chars: int = 500_000) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _limit(_extract_pdf(path), max_chars)
    if ext == ".docx":
        return _limit(_extract_docx(path), max_chars)
    if ext == ".pptx":
        return _limit(_extract_pptx(path), max_chars)
    if ext in {".xlsx", ".xls"}:
        return _limit(_extract_xlsx(path), max_chars)
    if ext in {".html", ".htm"}:
        return _limit(_extract_html(path), max_chars)
    if ext in TEXT_EXTENSIONS:
        return _limit(_read_text(path), max_chars)

    mime_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    size = path.stat().st_size if path.exists() else 0
    return f"[Unsupported binary source: {path.name}, mime={mime_type}, size={size} bytes]"


def _limit(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n\n[Truncated by backend extraction limit]\n"


def _read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="utf-8", errors="replace")


def _extract_pdf(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
    except Exception as exc:
        raise ExtractionError("PDF extraction requires PyMuPDF. Install backend requirements.") from exc

    chunks: list[str] = []
    with fitz.open(path) as doc:
        for index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                chunks.append(f"## Page {index}\n\n{text}")
    return "\n\n".join(chunks).strip()


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except Exception as exc:
        raise ExtractionError("DOCX extraction requires python-docx.") from exc

    document = Document(path)
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts).strip()


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise ExtractionError("PPTX extraction requires python-pptx.") from exc

    presentation = Presentation(path)
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"## Slide {index}\n\n" + "\n\n".join(texts))
    return "\n\n".join(slides).strip()


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        raise ExtractionError("XLSX extraction requires openpyxl.") from exc

    workbook = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                parts.append("\t".join(values).rstrip())
    return "\n".join(parts).strip()


def _extract_html(path: Path) -> str:
    raw = _read_text(path)
    try:
        from bs4 import BeautifulSoup
    except Exception:
        without_tags = re.sub(r"<[^>]+>", " ", raw)
        return html.unescape(re.sub(r"\s+", " ", without_tags)).strip()

    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return soup.get_text("\n", strip=True)


def slugify(value: str) -> str:
    base = Path(value).stem or "source"
    slug = re.sub(r"[^\w\u4e00-\u9fff.-]+", "-", base, flags=re.UNICODE).strip("-._")
    return slug or "source"


def markdown_table(rows: Iterable[Iterable[str]]) -> str:
    prepared = [[str(cell).replace("|", "\\|") for cell in row] for row in rows]
    if not prepared:
        return ""
    header = prepared[0]
    body = prepared[1:]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join("---" for _ in header) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def compact_json(value: object) -> str:
    return json.dumps(value, ensure_ascii=False, separators=(",", ":"))
